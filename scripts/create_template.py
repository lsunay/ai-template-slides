#!/usr/bin/env python3
"""
Script to create basic PowerPoint template files for the project.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_template(template_name: str, output_dir: str):
    """Create a basic template file."""
    output_path = Path(output_dir) / template_name / "template.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation()
    
    # Create title slide layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    if title:
        title.text = "Title"
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
    
    if subtitle:
        subtitle.text = "Subtitle"
        subtitle_frame = subtitle.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(28)
    
    # Create content slide layout
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    if title:
        title.text = "Slide Title"
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
    
    if content:
        content_frame = content.text_frame
        content_frame.clear()
        
        for i in range(3):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"Bullet point {i + 1}"
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    
    # Save the template
    prs.save(output_path)
    print(f"Template created: {output_path}")


if __name__ == "__main__":
    templates = ["academic", "pitch_deck", "sales"]
    output_dir = "backend/app/templates"
    
    for template in templates:
        create_template(template, output_dir)
    
    print("All templates created successfully!")