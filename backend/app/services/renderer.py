import uuid
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..core.config import settings


class PresentationRenderer:
    """Renderer for creating PowerPoint presentations from slide outlines."""
    
    def __init__(self):
        self.output_dir = Path(settings.output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    def render(self, template_path: Path, outline: dict, presentation_title: str) -> Path:
        """
        Render a presentation from template and outline.
        
        Args:
            template_path: Path to template .pptx file
            outline: Slide outline with titles and bullets
            presentation_title: Title for the presentation
        
        Returns:
            Path to generated presentation file
        """
        # Load template
        prs = Presentation(template_path)
        
        # Remove default slides
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]
        
        # Add title slide
        self._add_title_slide(prs, presentation_title)
        
        # Add content slides
        for i, title in enumerate(outline["titles"]):
            bullets = outline["bullets"][i] if i < len(outline["bullets"]) else []
            self._add_content_slide(prs, title, bullets)
        
        # Save presentation
        output_filename = f"{uuid.uuid4()}.pptx"
        output_path = self.output_dir / output_filename
        prs.save(output_path)
        
        return output_path
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add a title slide to the presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            self._format_title_shape(title_shape)
    
    def _add_content_slide(self, prs: Presentation, title: str, bullets: list):
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            self._format_title_shape(title_shape)
        
        # Set content (bullets)
        if bullets:
            content_shape = slide.placeholders[1]
            if content_shape:
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()
                    
                    paragraph.text = bullet
                    paragraph.level = 0
                    paragraph.font.size = Pt(18)
                    paragraph.font.name = 'Calibri'
                    paragraph.space_before = Pt(6)
                    paragraph.space_after = Pt(6)
    
    def _format_title_shape(self, shape):
        """Format title shape with consistent styling."""
        text_frame = shape.text_frame
        text_frame.clear()
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = shape.text
        paragraph.font.size = Pt(44)
        paragraph.font.name = 'Calibri'
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = PP_ALIGN.LEFT