import base64
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .parser import PresentationStructure, SlideContent


class PresentationRenderer:
    """Renderer for creating PowerPoint presentations from presentation structures."""
    
    def __init__(self):
        self.presentation: Optional[Presentation] = None
    
    def render_presentation(self, structure: PresentationStructure) -> Presentation:
        """Render a presentation structure into a PowerPoint presentation."""
        self.presentation = Presentation()
        
        # Set presentation metadata
        self.presentation.core_properties.title = structure.title
        if structure.subtitle:
            self.presentation.core_properties.comments = structure.subtitle
        
        # Add title slide
        self._add_title_slide(structure)
        
        # Add content slides
        for slide_content in structure.slides:
            self._add_content_slide(slide_content)
        
        return self.presentation
    
    def _add_title_slide(self, structure: PresentationStructure):
        """Add a title slide to the presentation."""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = structure.title
            self._format_title_shape(title_shape)
        
        # Set subtitle if available
        if structure.subtitle:
            subtitle_shape = slide.placeholders[1]
            if subtitle_shape:
                subtitle_shape.text = structure.subtitle
                self._format_subtitle_shape(subtitle_shape)
    
    def _add_content_slide(self, slide_content: SlideContent):
        """Add a content slide to the presentation."""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_content.title
            self._format_title_shape(title_shape)
        
        # Set content
        content_shape = slide.placeholders[1]
        if content_shape:
            # Split content by newlines and add as bullet points
            content_lines = slide_content.content.strip().split('\n')
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default paragraph
            
            for i, line in enumerate(content_lines):
                line = line.strip()
                if not line:
                    continue
                
                # Remove bullet characters if present
                line = line.lstrip('•-*').strip()
                
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                
                paragraph.text = line
                paragraph.level = 0
                paragraph.font.size = Pt(18)
                paragraph.font.name = 'Calibri'
                paragraph.space_before = Pt(6)
                paragraph.space_after = Pt(6)
        
        # Add speaker notes if available
        if slide_content.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_content.notes
    
    def _format_title_shape(self, shape):
        """Format a title shape with consistent styling."""
        text_frame = shape.text_frame
        text_frame.clear()
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = shape.text
        paragraph.font.size = Pt(44)
        paragraph.font.name = 'Calibri'
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
        paragraph.alignment = PP_ALIGN.LEFT
    
    def _format_subtitle_shape(self, shape):
        """Format a subtitle shape with consistent styling."""
        text_frame = shape.text_frame
        text_frame.clear()
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = shape.text
        paragraph.font.size = Pt(28)
        paragraph.font.name = 'Calibri'
        paragraph.font.color.rgb = RGBColor(89, 89, 89)  # Dark gray
        paragraph.alignment = PP_ALIGN.LEFT
    
    def save_to_file(self, file_path: str):
        """Save the presentation to a file."""
        if not self.presentation:
            raise ValueError("No presentation to save. Call render_presentation first.")
        
        self.presentation.save(file_path)
    
    def to_base64(self) -> str:
        """Convert the presentation to base64 string."""
        if not self.presentation:
            raise ValueError("No presentation to convert. Call render_presentation first.")
        
        buffer = BytesIO()
        self.presentation.save(buffer)
        buffer.seek(0)
        
        return base64.b64encode(buffer.read()).decode('utf-8')
    
    def from_base64(self, base64_string: str) -> Presentation:
        """Load a presentation from base64 string."""
        buffer = BytesIO(base64.b64decode(base64_string))
        self.presentation = Presentation(buffer)
        return self.presentation